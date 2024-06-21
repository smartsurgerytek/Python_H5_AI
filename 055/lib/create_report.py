import cv2
import numpy as np
import os
import pandas as pd
import numpy
from openpyxl import Workbook, load_workbook
from openpyxl.styles import PatternFill,Font
import pptx
from pptx.util import Inches,Pt
from pptx import Presentation
import time
"""
Create AI model predicted images report, content record IoU, misjudgment and omission information.
Notice: If image resolution is higher, the running time will be longer.
"""

class CreateAccuracyReport(object):
    """
    Inputs-
    trueLabelFolderPath: The folder path of transformed label data(ground truth, type is .png).
    predRootFolderPath: The root folder path of all prediction image folder(name format should be pred-xx, xx is number of model training epoch).
    reportTitleName: The title name used for Excel file name.
    colorDict: A dictionary that KEY is the number corresponding to the category, VALUE is the BGR list corresponding to the category.
    """
    # def __init__(self, trueLabelFolderPath, predRootFolderPath, reportTitleName, colorDict):
    def __init__(self, trueLabelFolderPath, predRootFolderPath, reportTitleName, config):#--new20220810-yita
        self.trueLabelFolderPath = trueLabelFolderPath
        self.predRootFolderPath = predRootFolderPath
        self.reportTitleName = reportTitleName
        
    #--------------------------new20220810-yita----------------  
        self.config=config
        self.sequence = list(map(lambda x: chr(x), range(ord('A'), ord('Z') + 1)))
        colorDict = {} # 格式為字典, KEY為類別對應到模型上的數值, VALUE為BGR顏色的清單(e.g [B, G, R])
        colorDict[0] = [0,0,0]
        for idx, color in enumerate(config["color_list"]):
            colorDict[idx+1] = color
        # print(color_dict)   

        detect_label_dict= {}
        detect_label_dict[0] = "Background"
        for idx, detect_label_list in enumerate(config["detect_label_list"]):
            detect_label_dict[idx+1] = detect_label_list
        # print(detect_label_dict)
        self.detect_label_dict=detect_label_dict
    #--------------------------new20220810-yita---------------- 
        self.compare_img_folder_path1=config["save_compare_img_folder_path"]
        self.colorDict = colorDict
        # print(self.colorDict)
        self.numClass = len(colorDict)
        # print(self.numClass)
        self.predRootFolderList = os.listdir(self.predRootFolderPath)
        # print(self.predRootFolderList)
        self.predRootFolderListNEW=[]
        for i in self.predRootFolderList:
            try:
                if i.split('.')[1]=='xlsx' or i.split('.')[1]=='pptx' or i.split('.')[1]=='ipynb_checkpoints':
                    print('Delete_list_name:'+i)
                else :
                    self.predRootFolderListNEW.append(i)
            except:
                self.predRootFolderListNEW.append(i)
        # print(self.predRootFolderListNEW)
        self.predRootFolderList=self.predRootFolderListNEW
        self.epochList = [folderName.split('-')[1] for folderName in self.predRootFolderList if '.ipynb_checkpoints' not in folderName]
        self.saveNameList = [self.reportTitleName + '-' + epoch + '.xlsx' for epoch in self.epochList]

        self.setColumnNames()#--new20220810-yita
        
        self.createDataframe()
        
        self.run()
    #--------------------------new20220810-yita----------------    
    # 
    def RGBToHex(self,r, g, b):
        return '#%02X%02X%02X' % (r, g, b)

    def blackwhite(self,my_hex):
        """Returns complementary RGB color

        Example:
        >>>complementaryColor('FFFFFF')
        '000000'
        """
        if my_hex[0] == '#':
            my_hex = my_hex[1:]
        rgb = (my_hex[0:2], my_hex[2:4], my_hex[4:6])
        comp = ['%02X' % (255 - int(a, 16)) for a in rgb]
        return ''.join(comp)

    def correct_label_color(self,saveExcelPath):    
    #依原始 的名稱定義的資料: dfSummary
        correct_label_color_dict= {}
        correct_label_color_dict[0] = "Background",[0,0,0]
        detect_label_list=self.config["detect_label_list"]
        print(detect_label_list[0])
        list_num=0
        for idx, color in enumerate(self.config["color_list"]):
            name=detect_label_list[list_num]
            correct_label_color_dict[idx+1] =name, color
            list_num+=1

        wb = load_workbook(saveExcelPath) #開啟活頁簿
        # ws = wb.active
        ws = wb.get_sheet_by_name(u"Summary")
        # ws1.sheet_properties.tabColor = "1072BA"  # 设置 sheet 标签背景色
        x,y=self.dfSummary.shape
        OtherData0 = ['']

        for i in range(y-numpy.size(OtherData0)):
            OtherData0.append('') 
        ws.append(OtherData0)


        OtherData2 = ['色彩名稱表','HEX','B','G','R','Color']
        # print(self.dfSummary.shape)
        OtherData2NameLength=numpy.size(OtherData2)
        for i in range(y-numpy.size(OtherData2)):
            OtherData2.append('') 
        ws.append((OtherData2))
        
        r=x+2+2
        for i in range(1,OtherData2NameLength+1):           
            row = ws.cell(r-1, i) 
            row.fill = PatternFill("solid", start_color="FFFF00")

        for num, key in correct_label_color_dict.items():    
            ws.cell(r, 1).value = key[0]
            ws.cell(r, 2).value = self.RGBToHex(key[1][2],key[1][1],key[1][0])
            ws.cell(r, 3).value = key[1][0]
            ws.cell(r, 4).value = key[1][1]
            ws.cell(r, 5).value = key[1][2]        
            ws.cell(r, 6).value = key[0]
            ws.cell(r, 6).font =  Font(color=self.blackwhite(ws.cell(r, 2).value))
            ws.cell(r, 6).fill = PatternFill(start_color=self.RGBToHex(key[1][2],key[1][1],key[1][0]).replace('#', ''), fill_type="solid") #used hex code for red color
            r=r+1
        wb.save(saveExcelPath)           #存檔
        wb.close()  # close the workbook

    #     
    def SummaryNames(self,rowCount,file):
        # rowCount = 0
        colsCount= 1
        fileName = os.path.splitext(file)[0]
        excelalgo = [fileName]
        # print("self.numClass:"+str(self.numClass))
        # Count_Details!B2   self.initial_count

        for i in range(self.numClass+1):
            # excelalgofu='=IF(IoU_Details!'+self.ten2TwentySix(colsCount)+str(rowCount+2)+'>0.5, 1, IF(IoU_Details!'+self.ten2TwentySix(colsCount)+str(rowCount+2)+'="", 1, 0))'
            # excelalgofu='=IF(IoU_Details!'+self.ten2TwentySix(colsCount)+str(rowCount+2)+'> Count_Details!B'+str(self.initial_count+4) +', 1, IF(IoU_Details!'+self.ten2TwentySix(colsCount)+str(rowCount+2)+'="", 1, 0))'
            excelalgofu='=IF(IoU_Details!'+self.ten2TwentySix(colsCount)+str(rowCount+2)+'> $B$'+str(self.initial_count+10) +', 1, IF(IoU_Details!'+self.ten2TwentySix(colsCount)+str(rowCount+2)+'="", 1, 0))'
            excelalgo.append(excelalgofu)
            colsCount+=1
     
        return excelalgo




    def SummaryOtherData(self,rowCount):
        cwd = os.getcwd()
        Data_total = cwd.split('\C_pred')[0]
        # print(Data_total)
        Data_total_excel=Data_total+'\A_Data_preprocessing\Dentistry_ClassInformation.xlsx'
        
        # self.predRootFolderPath
        workbook = load_workbook(Data_total_excel)
        # 獲取表單
        sheet = workbook[u'Summary']
        # 讀取指定的單元格資料
        rows = sheet.rows
        for row in list(rows):  # 遍歷每行資料
            case = []   # 用於存放一行資料
            for c in row:  # 把每行的每個單元格的值取出來，存放到case裡
                case.append(c.value)
            # print(case)
        # case#這是最後一行

        x,y=self.dfSummary.shape

        OtherData0_0 = ['Label'] 
        for i in range(y-numpy.size(OtherData0_0)):
            OtherData0_0fu='=('+self.ten2TwentySix(i+1)+str(1)+')'
            OtherData0_0.append(OtherData0_0fu) 
        # print(numpy.size(OtherData))
        self.dfSummary.loc[rowCount] = OtherData0_0
        rowCount+=1




        OtherData0_1 = ['整體準確率表現平均值']
        for i in range(y-numpy.size(OtherData0_1)):
            OtherData0_1fu='=AVERAGE('+self.ten2TwentySix(i+1)+str(2)+':'+self.ten2TwentySix(i+1)+str(x+1)+')'
            OtherData0_1.append(OtherData0_1fu) 
        # print(numpy.size(OtherData))
        self.dfSummary.loc[rowCount] = OtherData0_1
        rowCount+=1

        OtherData0_2 = ['個別IOU準確率表現平均值']
        for i in range(y-numpy.size(OtherData0_2)):

            OtherData0_2fu='=IFERROR(AVERAGE(Count_Details!'+self.ten2TwentySix(i+1)+str(2)+':' +self.ten2TwentySix(i+1)+str(x+1)+'),"無參考價值或AI完全無法學習成效")'
            OtherData0_2.append(OtherData0_2fu) 
        # print(numpy.size(OtherData))
        self.dfSummary.loc[rowCount] = OtherData0_2
        rowCount+=1

        OtherData0_3 = ['','個別判定總數']
        for num in range(y-numpy.size(OtherData0_3)):
            if num==0:
                OtherData0_3.append(case[y-1])
            # elif num == 1:
            #     OtherData0_3.append(case[y-1])
            else:
                OtherData0_3.append(case[num])
        self.dfSummary.loc[rowCount] = OtherData0_3
        rowCount+=1

        OtherData0_4 = ['','漏判數']
        for i in range(y-numpy.size(OtherData0_4)):
            OtherData0_4fu='=COUNTIF(Count_Details!'+self.ten2TwentySix(i+2)+str(2)+':' +self.ten2TwentySix(i+2)+str(x+1)+', "Miss")'
            OtherData0_4.append(OtherData0_4fu) 
        self.dfSummary.loc[rowCount] = OtherData0_4
        rowCount+=1

        OtherData0_5 = ['','誤判數']
        for i in range(y-numpy.size(OtherData0_5)):
            OtherData0_5fu='=COUNTIF(Count_Details!'+self.ten2TwentySix(i+2)+str(2)+':' +self.ten2TwentySix(i+2)+str(x+1)+', "Misjudgment")'
            OtherData0_5.append(OtherData0_5fu) 
        self.dfSummary.loc[rowCount] = OtherData0_5
        rowCount+=1

        OtherData0 = ['']
        for i in range(y-numpy.size(OtherData0)):
            OtherData0.append('') 
        self.dfSummary.loc[rowCount] = OtherData0
        rowCount+=1

        OtherData1 = ['']
        for i in range(y-numpy.size(OtherData1)):
            OtherData1.append('') 

        self.dfSummary.loc[rowCount] = OtherData1
        rowCount+=1

        OtherData = ['IOU設定值', 0.5]
        for i in range(y-numpy.size(OtherData)):
            OtherData.append('') 
        self.dfSummary.loc[rowCount] = OtherData


#         self.dfIoUColNames = ['編號類別對應表', 'ClassName', '', '']
#         for i in range(self.numClass):
#             # self.dfIoUColNames.append('IoU ' + str(i))
#             self.dfIoUColNames.append(self.detect_label_dict[i])
#         self.dfCountColNames = self.dfIoUColNames.copy()
#         self.dfIoUColNames.append('Mean IoU')
    #--------------------------new20220810-yita----------------        
    def setColumnNames(self):
        self.dfIoUColNames = ['File name', 'Pixel accuracy']
        for i in range(self.numClass):
            # self.dfIoUColNames.append('IoU ' + str(i))
            self.dfIoUColNames.append(self.detect_label_dict[i])
        self.dfCountColNames = self.dfIoUColNames.copy()
        self.dfIoUColNames.append('Mean IoU')

    def createDataframe(self):
        # self.dfSummary = pd.DataFrame()
        # print(self.dfCountColNames)
        self.dfSummary = pd.DataFrame(columns=self.dfCountColNames)
        self.dfIoU = pd.DataFrame(columns=self.dfIoUColNames)
        self.dfCount = pd.DataFrame(columns=self.dfCountColNames)
        

    def readAndTransMask(self, labelPath, predFilePath):
        labelArray = cv2.imread(labelPath, cv2.IMREAD_GRAYSCALE)
        height, width = labelArray.shape
        predImg = cv2.imread(predFilePath)
        predArray = np.zeros([height, width])
        for i in np.arange(1, self.numClass):
            idx = np.where((predImg[:,:,0] == self.colorDict[i][0])
                       & (predImg[:,:,1] == self.colorDict[i][1])
                       & (predImg[:,:,2] == self.colorDict[i][2]))
            predArray[idx] = i
        return labelArray, predArray

    def index2position(self, idxPairs):
        rows = idxPairs[0]
        cols = idxPairs[1]
        return {(y, x) for (y, x) in zip(rows, cols)}

    def calIoUAndCountList(self, labelArray, predArray):
        iouList = []
        countList=[]
        # SummaryList=[]
        for i in range(self.numClass):
            # print(labelIdx)
            
            labelIdx = np.where(labelArray==i)
            # print(labelIdx)
            predIdx = np.where(predArray==i)
            if len(labelIdx[0]) == 0 and len(predIdx[0]) == 0:
                iouList.append(None)
                countList.append(None)
            elif len(labelIdx[0]) == 0:
                iouList.append(0)
                countList.append('Misjudgment')
            elif len(predIdx[0]) == 0:
                iouList.append(0)
                countList.append('Miss')
            else:
                labelPos = self.index2position(labelIdx)
                predPos = self.index2position(predIdx)
                union = labelPos.union(predPos)
                inter = labelPos.intersection(predPos)

                iouList.append(len(inter)/len(union))
                countList.append(len(inter)/len(union))
            # SummaryList='=IF(IoU_Details!C'+str(rowCount+2)+'>0.5, 1, IF(IoU_Details!C'+str(rowCount+2)+'="", 1, 0))'
        miou = np.nanmean(np.array(iouList, dtype=np.float64))
        iouList.append(miou)
        return iouList, countList

    def calImgAccuracy(self, predFolderPath, file):
        fileName = os.path.splitext(file)[0]
        tmpRowIoU = [fileName]
        tmpRowCount = [fileName]
        
        labelPath = os.path.join(self.trueLabelFolderPath, fileName + '.png')
        predFilePath = os.path.join(predFolderPath, fileName + '_pred.png')
        labelArray, predArray = self.readAndTransMask(labelPath, predFilePath)
        addArray = np.add(labelArray, np.multiply(predArray, -1))
        pixelAccuracy = len(np.where(addArray==0)[0]) / addArray.size        #整張圖預測的總區域的正確性(IOU)
        iouList, countList = self.calIoUAndCountList(labelArray, predArray)  #個別項目預測的正確性(IOU)
        
        tmpRowIoU.append(pixelAccuracy)      #把pixelAccuracy新增一行在tmpRowIoU
        tmpRowIoU.extend(iouList)            #把iouList接續在tmpRowIoU值的最後面
        tmpRowCount.append(pixelAccuracy)
        tmpRowCount.extend(countList)
        
        # print(iouList)
        # print(tmpRowCount)
        return tmpRowIoU, tmpRowCount
    

    def ten2TwentySix(self,num):
        L = []
        if num > 25:
            while True:
                d = int(num / 26)
                remainder = num % 26
                if d <= 25:
                    L.insert(0, self.sequence[remainder])
                    L.insert(0, self.sequence[d - 1])
                    break
                else:
                    L.insert(0, self.sequence[remainder])
                    num = d - 1
        else:
            L.append(self.sequence[num])

        return "".join(L)
    # 建立PPT提供快速瀏覽
    def compare_output_ppt(self,predFolderName,saveExcelPath):

        time.localtime()
        train_time=time.strftime("%Y-%m-%d", time.localtime())
        # pred_mask_folder_path=config["save_pred_img_folder_path"]
        pred_mask_folder_path=predFolderName
        # 'pred-' + model +'_'+ backbone +'_' +ckptname

        print(predFolderName)
        User_name=self.config["user_name"]
        print("pred_mask_folder_path:   "+pred_mask_folder_path)
        savePPtPath=saveExcelPath.split('.xlsx')
        print(savePPtPath)
        ppt_filename=savePPtPath[0]+'.pptx'
        print(ppt_filename)
        # prs = Presentation()
        #開啟新的簡報物件
        
        prs = Presentation()
        #建立簡報檔第一張頁面物件
        title_slide_layout = prs.slide_layouts[0] 
        #增加一張簡報
        slide = prs.slides.add_slide(title_slide_layout)
        #設定第一張簡報的標題 
        title = slide.shapes.title
        # title.text = 'pred-' + model +'_'+ backbone +'_' +ckptname
        title.text = pred_mask_folder_path
        #設定第一張簡報的副標題
        subtitle = slide.placeholders[1]
        subtitle.text = "作者： " + User_name +  "   " +train_time
        #將簡報物件存檔
        prs.save(ppt_filename)
        # ppt_filename = 'python_ppt_v1'
        # full_ppt_filename = '{}.{}'.format(ppt_filename,'pptx')
        ppt_file = pptx.Presentation(ppt_filename)


        cwd = os.getcwd()
        # print(cwd)
        AI_Path = cwd.split('\D_build_report')[0]
        # print(AI_Path)
        # total_folder_path_new = pred_mask_folder_path.split('/')

        # label_mask_folder_path= AI_Path + '/C_pred/compare_output/' + total_folder_path_new[len(total_folder_path_new)-1]
        #label_mask_folder_path= AI_Path + '/C_pred/compare_output/' + pred_mask_folder_path #停用 modify by charley
        label_mask_folder_path= self.compare_img_folder_path1 + '/' + pred_mask_folder_path #停用 add by charley 0831
        print(label_mask_folder_path)
        # label_mask_folder_path='F:/SurgeryAnalytics/AI_Cases/Dentistry/2_Linknet_efficientnetb7_0_5/C_pred/compare_output/pred-Linknet_efficientnetb7_n41'

        # label_mask_folder_path=pred_mask_folder_path
        pic_files = [fn for fn in os.listdir(label_mask_folder_path) if fn.endswith('.png')]
        # print(pic_files)
        file_name_list = [os.path.splitext(x)[0] for x in os.listdir(label_mask_folder_path)]

        file_name_list = [os.path.splitext(x)[0] for x in os.listdir(label_mask_folder_path)]

        # fnp=label_mask_folder_path+'/'+file_name_list[0]+'.png'


        for i in range(0,len(file_name_list),4):

            slide = ppt_file.slides.add_slide(ppt_file.slide_layouts[1])
            # title = slide.shapes.title
            # slide.shapes.title.top = 0
            # slide.shapes.title.left = 0
            # # slide.shapes.title.height = 1
            # title.text = file_name_list[i] + "~" + str(i+3)
            textbox = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(0.25), width=(5), height=Inches(1.5))  # left，top爲相對位置，width，height爲文本框大小
            textbox.text = "清單順序" + str(i+1) + "~" + str(i+4)

            img_path1=label_mask_folder_path+'/'+file_name_list[i]+'.png'
            pic = slide.shapes.add_picture(img_path1, left=Inches(0), top=Inches(1.5), height=Inches(1.5))

            # img_path2=label_mask_folder_path+'/'+file_name_list[i+1]+'.png'
            # pic = slide.shapes.add_picture(img_path2, left=Inches(0), top=Inches(3), height=Inches(1.5))

            # img_path3=label_mask_folder_path+'/'+file_name_list[i+2]+'.png'
            # pic = slide.shapes.add_picture(img_path3, left=Inches(0), top=Inches(4.5), height=Inches(1.5))

            # img_path4=label_mask_folder_path+'/'+file_name_list[i+3]+'.png'
            # pic = slide.shapes.add_picture(img_path4, left=Inches(0), top=Inches(6), height=Inches(1.5))   
            try:
                img_path4=label_mask_folder_path+'/'+file_name_list[i+3]+'.png'
                pic = slide.shapes.add_picture(img_path4, left=Inches(0), top=Inches(6), height=Inches(1.5)) 
                img_path3=label_mask_folder_path+'/'+file_name_list[i+2]+'.png'
                pic = slide.shapes.add_picture(img_path3, left=Inches(0), top=Inches(4.5), height=Inches(1.5))
                img_path2=label_mask_folder_path+'/'+file_name_list[i+1]+'.png'
                pic = slide.shapes.add_picture(img_path2, left=Inches(0), top=Inches(3), height=Inches(1.5))
            except:
                try:
                    img_path3=label_mask_folder_path+'/'+file_name_list[i+2]+'.png'
                    pic = slide.shapes.add_picture(img_path3, left=Inches(0), top=Inches(4.5), height=Inches(1.5))
                    img_path2=label_mask_folder_path+'/'+file_name_list[i+1]+'.png'
                    pic = slide.shapes.add_picture(img_path2, left=Inches(0), top=Inches(3), height=Inches(1.5))
                except:
                    try:
                        img_path2=label_mask_folder_path+'/'+file_name_list[i+1]+'.png'
                        pic = slide.shapes.add_picture(img_path2, left=Inches(0), top=Inches(3), height=Inches(1.5))
                    except:
                        textbox = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(6), width=(5), height=Inches(1.5))  # left，top爲相對位置，width，height爲文本框大小
                        textbox.text = "全部圖片已完成"
            
            if i+4>=len(file_name_list):
                textbox.text=textbox.text + "_報告資料以至結尾"
            
        ppt_file.save(ppt_filename)

    # print(ten2TwentySix(50))
    def run(self):
        for idx in range(len(self.predRootFolderList)):
            predFolderName = self.predRootFolderList[idx]
            saveExcelName = self.saveNameList[idx]
            # print(saveExcelName)
            predFolderPath = os.path.join(self.predRootFolderPath, predFolderName)
            saveExcelPath = os.path.join(self.predRootFolderPath, saveExcelName)
            # print(saveExcelPath)
            self.compare_output_ppt(predFolderName,saveExcelPath)
            rowCount = 0
            print("Building {} report ...  ".format(predFolderName), end = '')


            self.initial_count = 0
            # dir = "RandomDirectory"
            for path in os.listdir(self.trueLabelFolderPath):
                if os.path.isfile(os.path.join(self.trueLabelFolderPath, path)):
                    self.initial_count += 1
            print(self.initial_count)




            for file in os.listdir(self.trueLabelFolderPath):
                tmpRowIoU, tmpRowCount = self.calImgAccuracy(predFolderPath, file)
                excelalgo=self.SummaryNames(rowCount,file)
                self.dfIoU.loc[rowCount] = tmpRowIoU
                self.dfCount.loc[rowCount] = tmpRowCount
                self.dfSummary.loc[rowCount] = excelalgo
                rowCount += 1

            self.SummaryOtherData(rowCount)
            
            # columns = [['A', 'A', 'B', 'B', 'C'], ['a', 'b', 'c', 'd', 'e']]# 創建形狀為（10，5） 的DataFrame 並設置二級標題
            # demo_df = pd.DataFrame(np.arange(50).reshape(10, 5), columns=columns)
            # print(demo_df)
            # self.style_df = self.style_color(demo_df, {"A": '#1C1C1C', "B": '#00EEEE', "C": '#1A1A1A'})
            
            self.dfIoU.replace('None', np.nan, inplace=True)
            self.dfCount.replace('None', np.nan, inplace=True)
            with pd.ExcelWriter(saveExcelPath) as writer:
                self.dfSummary.to_excel(writer, sheet_name='Summary', index=False)
                self.dfIoU.to_excel(writer, sheet_name='IoU_Details', index=False)
                self.dfCount.to_excel(writer, sheet_name='Count_Details', index=False)
                # self.style_df.to_excel(writer, sheet_name='sheet_name')
            self.correct_label_color(saveExcelPath)
            print("Done!")

   
    
  